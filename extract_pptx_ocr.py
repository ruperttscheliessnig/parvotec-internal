#!/usr/bin/env python3
"""Extract text + OCR screenshots from PPTX slides.
Per slide: text boxes + embedded image OCR → merged output.
"""
from __future__ import annotations
import io, sys
from pathlib import Path
from datetime import datetime

from pptx import Presentation
from pptx.util import Pt
from PIL import Image
import pytesseract

PPTX = Path.home() / "workspace-active/parvotec/Machine learning for Rupert/ASGCT2026/ASGCT recap Georg.pptx"
OUT_DIR = Path.home() / "workspace-active/parvotec/Machine learning for Rupert/transcripts"
OUT_DIR.mkdir(exist_ok=True)
OUT_MD = OUT_DIR / "ASGCT_recap_Georg_full.md"
IMG_DIR = OUT_DIR / "slide_images"
IMG_DIR.mkdir(exist_ok=True)

pytesseract.pytesseract.tesseract_cmd = "/opt/homebrew/bin/tesseract"


def ocr_image(img: Image.Image) -> str:
    """Run Tesseract OCR on a PIL image."""
    # Upscale small images for better OCR
    w, h = img.size
    if w < 800:
        scale = 800 / w
        img = img.resize((int(w * scale), int(h * scale)), Image.LANCZOS)
    text = pytesseract.image_to_string(img, lang="eng", config="--psm 6")
    return text.strip()


def extract_slide(slide, slide_num: int) -> tuple[str, str]:
    """Return (text_content, ocr_content) for one slide."""
    text_parts = []
    ocr_parts = []

    # 1. Extract all text boxes
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                line = para.text.strip()
                if line:
                    text_parts.append(line)

    # 2. Extract embedded images and OCR them
    for shape in slide.shapes:
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            try:
                img_bytes = shape.image.blob
                img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
                # Save for reference
                img_path = IMG_DIR / f"slide_{slide_num:02d}_{shape.shape_id}.png"
                img.save(str(img_path))
                # OCR
                ocr_text = ocr_image(img)
                if ocr_text and len(ocr_text) > 20:
                    ocr_parts.append(f"[Screenshot OCR]\n{ocr_text}")
                    print(f"    img OCR: {len(ocr_text)} chars")
            except Exception as e:
                ocr_parts.append(f"[Screenshot OCR failed: {e}]")

        # Also handle grouped shapes containing images
        elif hasattr(shape, "shapes"):
            for sub in shape.shapes:
                if sub.shape_type == 13:
                    try:
                        img_bytes = sub.image.blob
                        img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
                        img_path = IMG_DIR / f"slide_{slide_num:02d}_{sub.shape_id}_grp.png"
                        img.save(str(img_path))
                        ocr_text = ocr_image(img)
                        if ocr_text and len(ocr_text) > 20:
                            ocr_parts.append(f"[Group Screenshot OCR]\n{ocr_text}")
                            print(f"    group img OCR: {len(ocr_text)} chars")
                    except Exception:
                        pass

    return "\n".join(text_parts), "\n\n".join(ocr_parts)


def main():
    print(f"Loading: {PPTX.name}")
    prs = Presentation(str(PPTX))
    n = len(prs.slides)
    print(f"Slides: {n}\n")

    lines = [
        f"# ASGCT 2026 Recap — Georg Feichtinger",
        f"**Source:** `{PPTX.name}`  ",
        f"**Extracted:** {datetime.now().strftime('%Y-%m-%d %H:%M')}  ",
        f"**Slides:** {n}",
        "",
        "---",
        "",
    ]

    for i, slide in enumerate(prs.slides, 1):
        print(f"Slide {i:02d}/{n} ...", end=" ", flush=True)
        text, ocr = extract_slide(slide, i)
        has_content = bool(text.strip()) or bool(ocr.strip())
        print(f"text={len(text)} ocr={len(ocr)}")

        lines.append(f"## Slide {i:02d}")
        if text.strip():
            lines.append("\n**Text:**\n")
            lines.append(text)
        if ocr.strip():
            lines.append("\n**OCR from screenshots:**\n")
            lines.append(ocr)
        if not has_content:
            lines.append("*(no extractable content)*")
        lines.append("\n---\n")

    content = "\n".join(lines)
    OUT_MD.write_text(content, encoding="utf-8")
    print(f"\nDone → {OUT_MD}")
    print(f"Total chars: {len(content)}")


if __name__ == "__main__":
    main()
