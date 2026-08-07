#!/usr/bin/env python3
"""Transcribe all ASGCT2026 MP4 videos and extract PPTX slide text.
Output: one combined analysis document per source file + a master index.
"""
from __future__ import annotations
import subprocess, sys, json
from datetime import datetime
from pathlib import Path

PARVOTEC = Path(__file__).parent / "Machine learning for Rupert"
ASGCT_DIR = PARVOTEC / "ASGCT2026"
OUT_DIR = PARVOTEC / "transcripts"
OUT_DIR.mkdir(exist_ok=True)

MASTER = OUT_DIR / "ASGCT2026_master.md"
PYTHON = sys.executable


def extract_audio(mp4: Path) -> Path:
    wav = OUT_DIR / (mp4.stem + ".wav")
    if wav.exists():
        print(f"  audio cached: {wav.name}")
        return wav
    print(f"  extracting audio from {mp4.name} ...", end=" ", flush=True)
    subprocess.run(
        ["ffmpeg", "-y", "-i", str(mp4), "-ar", "16000", "-ac", "1",
         "-c:a", "pcm_s16le", str(wav)],
        capture_output=True, check=True
    )
    print("OK")
    return wav


def transcribe(wav: Path) -> str:
    txt = OUT_DIR / (wav.stem + ".txt")
    if txt.exists():
        print(f"  transcript cached: {txt.name}")
        return txt.read_text()
    print(f"  transcribing {wav.name} (medium model, ~5-15 min) ...", flush=True)
    result = subprocess.run(
        [PYTHON, "-c",
         f"""
import whisper, json, sys
model = whisper.load_model("medium")
r = model.transcribe("{wav}", language="en", fp16=False,
                     initial_prompt="AAV capsid gene therapy machine learning deep learning protein engineering")
print(r["text"])
"""],
        capture_output=True, text=True, timeout=3600
    )
    text = result.stdout.strip()
    if not text:
        text = f"[Transcription failed: {result.stderr[:300]}]"
    txt.write_text(text)
    return text


def extract_pptx(pptx: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(pptx))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
        if texts:
            slides.append(f"### Slide {i}\n" + "\n".join(texts))
    return "\n\n".join(slides)


def build_master(entries: list[dict]) -> None:
    lines = [
        "# ASGCT 2026 — AI/ML AAV Engineering: Master Knowledge Document",
        f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}",
        f"Source: {ASGCT_DIR}",
        "",
        "---",
        "",
        "## Index",
    ]
    for e in entries:
        lines.append(f"- [{e['title']}](#{e['anchor']})")
    lines += ["", "---", ""]

    for e in entries:
        lines += [
            f"## {e['title']} {{#{e['anchor']}}}",
            f"**File:** `{e['file']}`  ",
            f"**Type:** {e['type']}",
            "",
        ]
        if e.get("pre_analysis"):
            lines += ["### Pre-analysis", e["pre_analysis"], ""]
        lines += ["### Full text / transcript", "", e["content"], "", "---", ""]

    MASTER.write_text("\n".join(lines), encoding="utf-8")
    print(f"\nMaster document written: {MASTER}")


# ── Pre-analysis stubs (filled from filenames — expand after reading) ─────────
PRE_ANALYSIS = {
    "AAV_Engineering_III": (
        "AAV Engineering III session. Expected: ML-guided capsid engineering, "
        "directed evolution surrogates, multi-trait optimisation approaches."
    ),
    "AAV_Engineering_IV": (
        "AAV Engineering IV session. Expected: continuation of engineering methods, "
        "possibly structure-guided or LLM-based capsid design."
    ),
    "AAV_Trafficking": (
        "AAV Trafficking session. Expected: receptor interactions, intracellular "
        "trafficking, ML models for tropism prediction."
    ),
    "Lir_AAV_LLM": (
        "Lir group — AAV + LLM. Expected: large language models applied directly "
        "to AAV capsid sequence design or functional prediction."
    ),
    "ShapeTX_AAV5engineering": (
        "ShapeTX — AAV5 engineering. Expected: proprietary ML-based platform for "
        "AAV5 capsid diversification and functional screening."
    ),
    "TuningReceptorInteractions_Caltech": (
        "Caltech — Tuning receptor interactions. Expected: computational/experimental "
        "approaches to engineering AAV-receptor binding specificity."
    ),
    "ASGCT recap Georg": (
        "Georg Feichtinger's personal ASGCT 2026 recap. Key findings, "
        "highlighted talks, relevance for Parvotec ML project."
    ),
}

# ── Main ─────────────────────────────────────────────────────────────────────
entries = []

# 1. Videos
mp4_files = sorted(ASGCT_DIR.glob("*.mp4"))
print(f"\nFound {len(mp4_files)} videos.\n")
for mp4 in mp4_files:
    print(f"=== {mp4.name} ===")
    wav = extract_audio(mp4)
    transcript = transcribe(wav)
    anchor = mp4.stem.lower().replace(" ", "-").replace("_", "-")
    entries.append({
        "title": mp4.stem.replace("_", " "),
        "anchor": anchor,
        "file": mp4.name,
        "type": "Video transcript (Whisper medium)",
        "pre_analysis": PRE_ANALYSIS.get(mp4.stem, ""),
        "content": transcript,
    })
    print(f"  transcript: {len(transcript)} chars\n")

# 2. PPTX
pptx_files = sorted(ASGCT_DIR.glob("*.pptx"))
print(f"Found {len(pptx_files)} PPTX files.\n")
for pptx in pptx_files:
    print(f"=== {pptx.name} ===")
    text = extract_pptx(pptx)
    anchor = pptx.stem.lower().replace(" ", "-").replace("_", "-")
    entries.append({
        "title": pptx.stem,
        "anchor": anchor,
        "file": pptx.name,
        "type": "PowerPoint slide text",
        "pre_analysis": PRE_ANALYSIS.get(pptx.stem, ""),
        "content": text,
    })
    print(f"  slides extracted: {text.count('### Slide')} | {len(text)} chars\n")

build_master(entries)
